"""PPTX document parser using python-pptx."""
from pathlib import Path
from typing import Union, List
from uuid import UUID
from pptx import Presentation

from app.services.parsing.base import BaseParser
from app.services.parsing.models import ParsedDocument, ParsedSection
from app.services.parsing.normalizer import TextNormalizer
from app.core.logging import get_logger

logger = get_logger(__name__)


class PPTXParser(BaseParser):
    """Parser for PPTX documents using python-pptx."""
    
    def __init__(self):
        """Initialize PPTX parser."""
        self.normalizer = TextNormalizer()
    
    def supports_format(self, file_extension: str) -> bool:
        """Check if this parser supports PPTX format."""
        return file_extension.lower() in ['pptx']
    
    def _extract_text_from_shape(self, shape) -> str:
        """
        Extract text from a shape, preserving bullet hierarchy.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            str: Extracted text
        """
        if not hasattr(shape, "text_frame"):
            return ""
        
        text_parts = []
        
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Get indentation level (0-based)
                level = paragraph.level
                
                # Add indentation for nested bullets
                indent = "  " * level
                
                # Preserve bullet points
                text_parts.append(f"{indent}{text}")
        
        return '\n'.join(text_parts)
    
    async def parse(
        self,
        file_path: Union[str, Path],
        document_id: UUID,
    ) -> ParsedDocument:
        """
        Parse PPTX document slide by slide.
        
        Features:
        - Extract slide title
        - Extract text from shapes
        - Extract speaker notes
        - Preserve bullet hierarchy
        - Use slide number as page_number
        
        Args:
            file_path: Path to PPTX file
            document_id: Document UUID
            
        Returns:
            ParsedDocument: Parsed PPTX with sections (one per slide)
            
        Raises:
            Exception: If PPTX parsing fails
        """
        file_path = Path(file_path)
        
        logger.info(
            f"Parsing PPTX document",
            extra={
                "document_id": str(document_id),
                "file_path": str(file_path),
            }
        )
        
        try:
            # Open PPTX presentation
            prs = Presentation(file_path)
            
            sections: List[ParsedSection] = []
            total_slides = len(prs.slides)
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, start=1):
                content_parts = []
                slide_title = None
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    # Check if shape has title
                    if hasattr(shape, "name") and "Title" in shape.name:
                        if hasattr(shape, "text"):
                            slide_title = shape.text.strip()
                    
                    # Extract text from shape
                    text = self._extract_text_from_shape(shape)
                    if text:
                        content_parts.append(text)
                
                # Extract speaker notes if present
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        notes_text = notes_frame.text.strip()
                        content_parts.append(f"\n[Speaker Notes]\n{notes_text}")
                
                # Skip empty slides
                if not content_parts and not slide_title:
                    logger.debug(
                        f"Skipping empty slide {slide_num}",
                        extra={"document_id": str(document_id)}
                    )
                    continue
                
                # Combine all content
                content = '\n\n'.join(content_parts)
                normalized_content = self.normalizer.normalize(content)
                
                # Use slide number as section title if no title found
                if not slide_title:
                    slide_title = f"Slide {slide_num}"
                else:
                    slide_title = self.normalizer.normalize_section_title(slide_title)
                
                # Create section for this slide
                section = ParsedSection(
                    section_title=slide_title,
                    content=normalized_content,
                    page_number=slide_num,
                    metadata={
                        "slide_number": slide_num,
                        "shape_count": len(slide.shapes),
                        "has_notes": slide.has_notes_slide,
                    }
                )
                
                sections.append(section)
            
            # Create parsed document
            parsed_doc = ParsedDocument(
                document_id=document_id,
                sections=sections,
                total_pages=total_slides,
                metadata={
                    "parser": "pptx",
                    "total_slides": total_slides,
                    "non_empty_slides": len(sections),
                }
            )
            
            # Log parsing statistics
            self._log_parsing_stats(
                document_id=document_id,
                sections_count=len(sections),
                total_chars=parsed_doc.get_total_content_length(),
                total_pages=total_slides,
            )
            
            return parsed_doc
            
        except Exception as e:
            logger.error(
                f"Failed to parse PPTX: {str(e)}",
                extra={
                    "document_id": str(document_id),
                    "file_path": str(file_path),
                },
                exc_info=True,
            )
            raise Exception(f"PPTX parsing failed: {str(e)}")
